"""RichFileTool -- Read docx/pdf/xlsx/pptx files with Python parsing libraries."""

# 试验243林晨
from pathlib import Path
from typing import Type

from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field


class RichFileInput(BaseModel):
    file_path: str = Field(
        description="Absolute or relative path to the file to read"
    )


class RichFileReadTool(BaseTool):
    """Read rich format files: docx, pdf, xlsx, pptx."""
    name: str = "read_rich_file"
    description: str = ""
    args_schema: Type[BaseModel] = RichFileInput
    root_dir: str = ""

    def _run(self, file_path: str) -> str:
        try:
            root = Path(self.root_dir)
            normalized = file_path.replace("\\", "/").lstrip("./")
            full_path = (root / normalized).resolve() if not Path(file_path).is_absolute() else Path(file_path).resolve()

            # Sandbox check
            if self.root_dir and not str(full_path).startswith(str(root.resolve())):
                # Also allow uploads dir
                upload_dir = root / "uploads"
                if not str(full_path).startswith(str(upload_dir.resolve())):
                    return f"Access denied: path escapes allowed directory"

            if not full_path.exists():
                return f"File not found: {file_path}"

            if not full_path.is_file():
                return f"Not a file: {file_path}"

            ext = full_path.suffix.lower()

            if ext == ".docx":
                return self._read_docx(full_path)
            elif ext == ".pdf":
                return self._read_pdf(full_path)
            elif ext in (".xlsx", ".xls"):
                return self._read_xlsx(full_path)
            elif ext == ".pptx":
                return self._read_pptx(full_path)
            else:
                # Fallback to plain text
                try:
                    content = full_path.read_text(encoding="utf-8")
                    if len(content) > 20000:
                        content = content[:20000] + "\n...[truncated]"
                    return content
                except UnicodeDecodeError:
                    return f"Cannot read binary file: {ext}"

        except ImportError as e:
            return f"Missing library for {ext}: {e}. Install with pip."
        except Exception as e:
            return f"Error reading file: {str(e)}"

    def _read_docx(self, path: Path) -> str:
        """Read .docx file using python-docx."""
        from docx import Document
        doc = Document(str(path))
        lines = []
        for para in doc.paragraphs:
            if para.text.strip():
                lines.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append(" | ".join(cells))
        content = "\n".join(lines)
        if len(content) > 20000:
            content = content[:20000] + "\n...[truncated]"
        return content

    def _read_pdf(self, path: Path) -> str:
        """Read .pdf file using PyPDF2."""
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        lines = []
        for i, page in enumerate(reader.pages[:50]):  # Max 50 pages
            text = page.extract_text()
            if text:
                lines.append(f"--- Page {i+1} ---\n{text}")
        content = "\n\n".join(lines)
        if len(content) > 20000:
            content = content[:20000] + "\n...[truncated]"
        return content

    def _read_xlsx(self, path: Path) -> str:
        """Read .xlsx file using openpyxl."""
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames[:5]:  # Max 5 sheets
            ws = wb[sheet_name]
            lines.append(f"=== Sheet: {sheet_name} ===")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 100:  # Max 100 rows per sheet
                    lines.append("...[more rows truncated]")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                lines.append(" | ".join(cells))
        wb.close()
        content = "\n".join(lines)
        if len(content) > 20000:
            content = content[:20000] + "\n...[truncated]"
        return content

    def _read_pptx(self, path: Path) -> str:
        """Read .pptx file using python-pptx."""
        from pptx import Presentation
        prs = Presentation(str(path))
        lines = []
        for i, slide in enumerate(prs.slides[:30]):  # Max 30 slides
            lines.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            lines.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        lines.append(" | ".join(cells))
        content = "\n".join(lines)
        if len(content) > 20000:
            content = content[:20000] + "\n...[truncated]"
        return content


def create_rich_file_tool(base_dir: Path) -> RichFileReadTool:
    """Create rich file reading tool."""
    tool = RichFileReadTool(root_dir=str(base_dir))
    tool.description = (
        "[Rich File Reader] Read rich format files: .docx, .pdf, .xlsx, .pptx. "
        "Supports extracting text, tables, and slides from these formats. "
        "Also falls back to plain text reading for .txt/.md/.py etc. "
        "Input: file path (relative to project root or absolute). "
        "Features: docx paragraphs+tables, pdf text extraction, xlsx multi-sheet, pptx slides+tables."
    )
    return tool
